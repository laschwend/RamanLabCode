# -*- coding: utf-8 -*-
"""
Created on Tue Jun 17 21:01:20 2025

@author: laSch

purpose: to go through a directory with a specific organization and put it into a powerpoint presentation bc I'm to lazy to copy paste stuff'
"""

import os 
import math
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from PIL import Image

"""
user defined variables
"""

rootDirectoryName = "C:/Users/laSch/MIT Dropbox/Raman Lab/Laura Schwendeman/NMJ Muscle Only Experiments"

fileTypesNames = ('.png', '.jpg', '.jpeg', '.bmp', '.tif')


"""
Actual code to make the powerpoint
"""
for (root, dirs, files) in os.walk(rootDirectoryName, topdown=True):
    
    print( print("Directory path: %s"%root))
    
    rel_path = os.path.relpath(root, rootDirectoryName)
    if rel_path == ".":
        depth = -1
    else:
        depth = rel_path.count(os.sep)
        
    #make new presentation
    if depth == 0:
        output_pptx = root + '/output_presentation.pptx'
        prs = Presentation()
        blank_slide_layout = prs.slide_layouts[6]  # Blank slide layout

        # Get slide dimensions
        slide_width = prs.slide_width
        slide_height = prs.slide_height

        # Padding (optional)
        margin_inch = 0.5
        max_width = slide_width - Inches(2 * margin_inch)
        max_height = slide_height - Inches(2 * margin_inch + 0.5)  # leave space for label at top

        #figure out how many directories to go through before saving
        numFoldersInExp = len([x for x in dirs])
      
        #set counter to 1
        print("Number of Folders in Exp:", numFoldersInExp)
        FolderCounter = 1
        
      
    #add a slide and put pictures in it   
    elif depth == 1:
        #don't do anything if there aren't any pictures in this folder bc then code breaks
        if len([x for x in files if x.lower().endswith(fileTypesNames)]) >0:
            print("here doing some stuff")
            #make a new slide
            slide = prs.slides.add_slide(blank_slide_layout)
            
            #get the name of the directory and make it the title
            dirName = os.path.basename(os.path.normpath(root))
            
            textbox = slide.shapes.add_textbox(Inches(0.5), Inches(0.1), Inches(9), Inches(0.5))
            text_frame = textbox.text_frame
            p = text_frame.paragraphs[0]
            p.text = dirName
            p.font.size = Pt(18)
            p.font.bold = True
            p.font.color.rgb = RGBColor(0, 0, 0)
            p.alignment = PP_ALIGN.CENTER
            
            #get the number of pictures
            numPictures = len([x for x in files if x.lower().endswith(fileTypesNames)])
            
            #do some math to figure out how to space everything in an appropriate grid
            gridSize = math.ceil((math.sqrt(numPictures)))
            
            x_spacing = math.floor(max_width/gridSize)
            y_spacing = math.floor(max_height/gridSize)
            
            
            #go through each file and put it in the grid and label it
            x_loc = 0
            y_loc = 0
            count = 1
            for filename in sorted(os.listdir(root)):
                if filename.lower().endswith(fileTypesNames):
                    print("y_loc", y_loc)
                    print("x_loc", x_loc)
                    print("gridSize", gridSize)
                    img_path = os.path.join(root, filename)
    
                     # Open image to get its size
                    with Image.open(img_path) as img:
                        img_width, img_height = img.size
                        img_ratio = img_width / img_height
                        box_ratio = x_spacing / y_spacing
    
                        # Scale proportionally
                        if img_ratio > box_ratio:
                            display_width = x_spacing
                            display_height = x_spacing / img_ratio
                        else:
                            display_height = y_spacing
                            display_width = y_spacing * img_ratio
                        
                    #figure out location of the image
                    left = margin_inch + display_width*x_loc
                    
                    top = Inches(0.5) + (y_spacing*.95)*y_loc + y_spacing*.05
                    
                    #now put the picture in the slide
                    pic = slide.shapes.add_picture(img_path, left, top, width=display_width, height=display_height)
                    
                    #now put text under the picture
                    left_t = left
                    top_t = top + display_height
                    width_t = display_width
                    height_t = display_height*.05
                    textbox = slide.shapes.add_textbox(left_t, top_t, width_t, height_t)
                    text_frame = textbox.text_frame
                    p = text_frame.paragraphs[0]
                    p.text = filename
                    p.font.size = Pt(8)
                    p.alignment = PP_ALIGN.CENTER
                    
                    count = count + 1
                    #now update next image location
                    if count <= gridSize:
                        #just move in x
                        x_loc = x_loc+1
                    else:
                        #reset x and add to y
                        x_loc = 0 
                        y_loc = y_loc + 1
                        count = 1
                    
        #now update the folder count processed and if done, save the presentation
        print(FolderCounter)
        if FolderCounter == numFoldersInExp:
            prs.save(output_pptx)
            print(f"Presentation saved to: {output_pptx}")
        else:
            FolderCounter = FolderCounter + 1 
                    
                
                
                
                
                
                
                
                
                
        
        
                             
                             
                             
                             
        
        
        
    
    
    
    
    
    
    
