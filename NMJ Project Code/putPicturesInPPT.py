# -*- coding: utf-8 -*-
"""
Created on Mon May 19 12:04:15 2025

@author: laSch
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from PIL import Image
import os

# --- Configuration ---
#image_folder = 'C:/Users/laSch/MIT Dropbox/Raman Lab/Angel Bu and Sonika Kohli/NMJ Stamping 3-Ch 5_15_25/Stamp Images'  # Change to your image folder
#image_folder = 'C:/Users/laSch/MIT Dropbox/Raman Lab/Laura Schwendeman/5_22_25 R6-7 stamps'
#image_folder = 'C:/Users/laSch/MIT Dropbox/Raman Lab/Laura Schwendeman/6_4_25 Opto G0 Seeding/Red Channel Pictures'
#image_folder = 'C:/Users/laSch/MIT Dropbox/Raman Lab/Laura Schwendeman/6_4_25 Opto Muscle Only Gels'
image_folder = 'C:/Users/laSch/MIT Dropbox/Raman Lab/Laura Schwendeman/6_8_2025 Opto D3 NMJ muscle only'
image_folder = 'C:/Users/laSch/MIT Dropbox/Raman Lab/Laura Schwendeman/7_25_25 leak test r1 - 11-8-345'
#image_folder = image_folder.replace(os.sep, '/')
image_folder = 'C:/Users/laSch/MIT Dropbox\Raman Lab\Laura Schwendeman/8_8_25 Leica Stamps from 8_7 before assembly'
image_folder = 'C:/Users/laSch/MIT Dropbox\Raman Lab\Laura Schwendeman/8_8_25 Leica Stamps from 8_7 before assembly'
image_folder = 'C:/Users/laSch/MIT Dropbox/Raman Lab/Laura Schwendeman/8_11_25 nmj GELMA casting results/DAY 0 Seeding Results processing Folder'
image_folder = 'C:/Users/laSch/MIT Dropbox/Raman Lab/Laura Schwendeman/9_18_25 Gel Longevity and Leak Studies/Gels D1'

output_pptx = image_folder + '/output_presentation.pptx'

# --- Create Presentation ---
prs = Presentation()
blank_slide_layout = prs.slide_layouts[6]  # Blank slide layout

# Get slide dimensions
slide_width = prs.slide_width
slide_height = prs.slide_height

# Padding (optional)
margin_inch = 0.5
max_width = slide_width - Inches(2 * margin_inch)
max_height = slide_height - Inches(2 * margin_inch + 0.5)  # leave space for label at top

# --- Loop through image files ---
for filename in sorted(os.listdir(image_folder)):
    if filename.lower().endswith(('.png', '.jpg', '.jpeg', '.bmp', '.tif')):
        slide = prs.slides.add_slide(blank_slide_layout)
        img_path = os.path.join(image_folder, filename)

         # Open image to get its size
        with Image.open(img_path) as img:
            img_width, img_height = img.size
            img_ratio = img_width / img_height
            box_ratio = max_width / max_height

            # Scale proportionally
            if img_ratio > box_ratio:
                display_width = max_width
                display_height = max_width / img_ratio
            else:
                display_height = max_height
                display_width = max_height * img_ratio

        # Center image
        left = (slide_width - display_width) / 2
        top = (slide_height - display_height) / 2 + Inches(0.2)
        
        # Add image to slide
        pic = slide.shapes.add_picture(img_path, left, top, width=display_width, height=display_height)  # Adjust size as needed

        # Add label (filename) as textbox
        textbox = slide.shapes.add_textbox(Inches(0.5), Inches(0.1), Inches(9), Inches(0.5))
        text_frame = textbox.text_frame
        p = text_frame.paragraphs[0]
        p.text = filename
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = RGBColor(0, 0, 0)
        p.alignment = PP_ALIGN.CENTER  # Corrected alignment

# --- Save presentation ---
prs.save(output_pptx)
print(f"Presentation saved to: {output_pptx}")